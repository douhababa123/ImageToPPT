from __future__ import annotations

import math
import re
from dataclasses import dataclass
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.config import settings
from app.services.ocr import OcrTextBox
from app.services.text_style import estimate_text_appearance
from app.services.text_rules import should_keep_as_background


PX_PER_INCH = 96


@dataclass(slots=True)
class ImagePlacement:
    left: int
    top: int
    width: int
    height: int
    scale_x: float
    scale_y: float


def build_pptx(slides: list[tuple[Path, Path, list[OcrTextBox]]], output_path: Path) -> Path:
    presentation = Presentation()

    first_image = Image.open(slides[0][0])
    first_width, first_height = first_image.size
    if settings.slide_mode == "source":
        presentation.slide_width = Inches(first_width / PX_PER_INCH)
        presentation.slide_height = Inches(first_height / PX_PER_INCH)
    else:
        presentation.slide_width = Inches(settings.slide_width_inches)
        presentation.slide_height = Inches(settings.slide_height_inches)

    blank_layout = presentation.slide_layouts[6]
    for original_image, cleaned_image, text_boxes in slides:
        with Image.open(original_image) as image:
            original_rgb = image.convert("RGB")
            original_rgb_array = np.asarray(original_rgb)
            image_width, image_height = image.size

        slide = presentation.slides.add_slide(blank_layout)
        placement = _calculate_image_placement(
            image_width=image_width,
            image_height=image_height,
            slide_width=presentation.slide_width,
            slide_height=presentation.slide_height,
        )
        slide.shapes.add_picture(
            str(cleaned_image),
            placement.left,
            placement.top,
            width=placement.width,
            height=placement.height,
        )

        for text_box in text_boxes:
            if should_keep_as_background(text_box, image_width, image_height):
                continue
            appearance = estimate_text_appearance(original_rgb_array, text_box)
            if appearance is None or appearance.light_text_on_colored_background:
                continue
            add_editable_text(slide, text_box, placement, original_rgb)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return output_path


def _calculate_image_placement(image_width: int, image_height: int, slide_width: int, slide_height: int) -> ImagePlacement:
    if settings.image_fit == "stretch":
        return ImagePlacement(
            left=0,
            top=0,
            width=slide_width,
            height=slide_height,
            scale_x=slide_width / image_width,
            scale_y=slide_height / image_height,
        )

    if settings.image_fit == "contain":
        scale = min(slide_width / image_width, slide_height / image_height)
    elif settings.image_fit == "cover":
        scale = max(slide_width / image_width, slide_height / image_height)
    else:
        raise ValueError(f"Unsupported IMAGE_FIT: {settings.image_fit}")

    rendered_width = int(image_width * scale)
    rendered_height = int(image_height * scale)
    return ImagePlacement(
        left=int((slide_width - rendered_width) / 2),
        top=int((slide_height - rendered_height) / 2),
        width=rendered_width,
        height=rendered_height,
        scale_x=scale,
        scale_y=scale,
    )


def add_editable_text(slide, text_box: OcrTextBox, placement: ImagePlacement, original_image: Image.Image) -> None:
    x1, y1, x2, y2 = text_box.bounds
    left = int(placement.left + x1 * placement.scale_x)
    top = int(placement.top + y1 * placement.scale_y)
    width = max(int((x2 - x1) * placement.scale_x), 1)
    height = max(int((y2 - y1) * placement.scale_y), 1)

    shape = slide.shapes.add_textbox(left, top, width, height)
    shape.rotation = _estimate_rotation(text_box)
    shape.fill.background()
    shape.line.fill.background()

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.word_wrap = False
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text_box.text

    estimated_font_size = _estimate_font_size(text_box.text, width, height)
    run.font.size = Pt(estimated_font_size)
    run.font.name = "Arial"
    run.font.bold = _should_use_bold(text_box.text, height)
    run.font.color.rgb = _sample_text_color(original_image, text_box)
    _mark_run_as_proofed(run)


def _mark_run_as_proofed(run) -> None:
    run_properties = run._r.get_or_add_rPr()
    run_properties.set("lang", "zh-CN")
    run_properties.set("dirty", "0")
    run_properties.set("smtClean", "1")
    run_properties.set("noProof", "1")


def _estimate_rotation(text_box: OcrTextBox) -> float:
    if len(text_box.box) < 2:
        return 0

    (x1, y1), (x2, y2) = text_box.box[0], text_box.box[1]
    angle = math.degrees(math.atan2(y2 - y1, x2 - x1))
    return angle if abs(angle) >= 2 else 0


def _estimate_font_size(text: str, width: int, height: int) -> float:
    height_pt = height / 12700
    width_pt = width / 12700
    weighted_char_count = sum(_character_width_weight(character) for character in text.strip())

    by_height = height_pt * 0.86
    by_width = width_pt / max(weighted_char_count, 1) * 1.06
    return max(5, min(96, by_height, by_width if weighted_char_count > 0 else by_height))


def _should_use_bold(text: str, height: int) -> bool:
    if height / 12700 >= 13:
        return True
    return bool(re.match(r"^\d+\.?\s+[A-Z]", text.strip()))


def _character_width_weight(character: str) -> float:
    if character.isspace():
        return 0.32
    if ord(character) <= 127:
        return 0.56
    return 1.0


def _sample_text_color(image: Image.Image, text_box: OcrTextBox) -> RGBColor:
    image_rgb = np.asarray(image.convert("RGB"))
    appearance = estimate_text_appearance(image_rgb, text_box)
    if appearance is None:
        return RGBColor(30, 30, 30)

    red, green, blue = appearance.text_rgb
    return RGBColor(red, green, blue)
